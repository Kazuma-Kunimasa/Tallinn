

import os
import time
from pathlib import Path

try:
    import PyPDF2
except ImportError:
    print("PyPDF2 is not installed. Please run 'pip install PyPDF2'")
    exit()

try:
    import pdfplumber
except ImportError:
    print("pdfplumber is not installed. Please run 'pip install pdfplumber'")
    exit()

try:
    from docx import Document
except ImportError:
    print("python-docx is not installed. Please run 'pip install python-docx'")
    exit()

try:
    import openpyxl
except ImportError:
    print("openpyxl is not installed. Please run 'pip install openpyxl'")
    exit()

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed. Please run 'pip install python-pptx'")
    exit()


class LocalFileConverter:
    def __init__(self, start_path):
        self.supported_extensions = {'.pdf', '.docx', '.xlsx', '.pptx'}
        self.conversion_stats = {
            'total': 0,
            'success': 0,
            'failed': 0,
            'errors': []
        }
        self.start_path = start_path

    def extract_text_from_pdf(self, file_path):
        text_content = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        text_content.append(f"=== ページ {page_num} ===\n{text}\n")
        except Exception as e:
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page_num, page in enumerate(reader.pages, 1):
                        text = page.extract_text()
                        if text:
                            text_content.append(f"=== ページ {page_num} ===\n{text}\n")
            except Exception as e2:
                raise Exception(f"PDF読み取りエラー (pdfplumber: {e}, PyPDF2: {e2})")
        return '\n'.join(text_content) if text_content else "テキストを抽出できませんでした"

    def extract_text_from_docx(self, file_path):
        try:
            doc = Document(file_path)
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_content.append(' | '.join(row_text))
            return '\n'.join(text_content) if text_content else "テキストを抽出できませんでした"
        except Exception as e:
            raise Exception(f"Word文書読み取りエラー: {e}")

    def extract_text_from_xlsx(self, file_path):
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== シート: {sheet_name} ===")
                for row in sheet.iter_rows():
                    row_data = [str(cell.value) for cell in row if cell.value is not None]
                    if row_data:
                        text_content.append(' | '.join(row_data))
                text_content.append("")
            return '\n'.join(text_content) if text_content else "テキストを抽出できませんでした"
        except Exception as e:
            raise Exception(f"Excelファイル読み取りエラー: {e}")

    def extract_text_from_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            text_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== スライド {slide_num} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                text_content.append("")
            return '\n'.join(text_content) if text_content else "テキストを抽出できませんでした"
        except Exception as e:
            raise Exception(f"PowerPointファイル読み取りエラー: {e}")

    def extract_text_from_file(self, file_path):
        file_ext = Path(file_path).suffix.lower()
        try:
            if file_ext == '.pdf':
                return self.extract_text_from_pdf(file_path)
            elif file_ext == '.docx':
                return self.extract_text_from_docx(file_path)
            elif file_ext == '.xlsx':
                return self.extract_text_from_xlsx(file_path)
            elif file_ext == '.pptx':
                return self.extract_text_from_pptx(file_path)
            else:
                raise Exception(f"サポートされていないファイル形式: {file_ext}")
        except Exception as e:
            error_msg = f"ファイル: {os.path.basename(file_path)}\nエラー: {str(e)}"
            self.conversion_stats['errors'].append(error_msg)
            raise e

    def convert_file_to_text(self, file_path):
        try:
            text_content = self.extract_text_from_file(file_path)
            header = f"=== 元ファイル: {os.path.basename(file_path)} ===\n"
            header += f"=== 変換日時: {time.strftime('%Y-%m-%d %H:%M:%S')} ===\n"
            header += f"=== ファイルパス: {file_path} ===\n\n"
            full_text = header + text_content
            
            text_file_path = Path(file_path).with_suffix('.txt')

            with open(text_file_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            
            self.conversion_stats['success'] += 1
            return True, text_file_path
        except Exception as e:
            self.conversion_stats['failed'] += 1
            return False, str(e)

    def process_folder(self):
        print("🚀 ファイル変換を開始...")
        print("=" * 50)

        target_files = []
        for root, dirs, files in os.walk(self.start_path):
            # venvディレクトリを除外
            if 'venv' in dirs:
                dirs.remove('venv')
            for file in files:
                if Path(file).suffix.lower() in self.supported_extensions:
                    target_files.append(os.path.join(root, file))

        if not target_files:
            print("❌ 変換対象のファイルが見つかりませんでした")
            return

        self.conversion_stats['total'] = len(target_files)
        print(f"📊 変換対象ファイル数: {len(target_files)}")
        print()

        for i, file_path in enumerate(target_files, 1):
            relative_path = os.path.relpath(file_path, self.start_path)
            print(f"[{i}/{len(target_files)}] 処理中: {relative_path}")
            
            success, result = self.convert_file_to_text(file_path)
            
            if success:
                print(f"✅ 完了: {os.path.basename(result)}")
            else:
                print(f"❌ エラー: {result}")
            print()

        self.show_conversion_summary()

    def show_conversion_summary(self):
        print("=" * 50)
        print("📊 変換結果サマリー")
        print(f"総ファイル数: {self.conversion_stats['total']}")
        print(f"成功: {self.conversion_stats['success']}")
        print(f"失敗: {self.conversion_stats['failed']}")
        
        if self.conversion_stats['errors']:
            print("\n❌ エラー詳細:")
            for error in self.conversion_stats['errors']:
                print(f"  {error}")
        
        success_rate = (self.conversion_stats['success'] / self.conversion_stats['total'] * 100) if self.conversion_stats['total'] > 0 else 0
        print(f"\n成功率: {success_rate:.1f}%")

if __name__ == '__main__':
    # スクリプトは現在の作業ディレクトリで実行されます。
    current_directory = os.getcwd()
    converter = LocalFileConverter(current_directory)
    converter.process_folder()

